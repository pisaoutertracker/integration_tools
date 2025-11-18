#!/usr/bin/env python3
"""
OMPO Report Editor - Edit PowerPoint presentation fields

This script modifies the yellow-highlighted fields in the OMPO Report PowerPoint presentation.
It can be used as a command-line tool or imported as a module.

Example usage:
    python edit_ompo_report.py input.pptx output.pptx --name "John Doe" --date "2025-11-18" --status "All systems operational"
"""

import argparse
import sys
from datetime import datetime
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import RGBColor


def find_and_replace_text_in_shape(shape, old_text, new_text):
    """
    Find and replace text in a shape, preserving formatting.
    
    Args:
        shape: The shape object to search
        old_text: Text to find
        new_text: Text to replace with
        
    Returns:
        bool: True if text was replaced, False otherwise
    """
    if not shape.has_text_frame:
        return False
    
    replaced = False
    text_frame = shape.text_frame
    
    # Check all paragraphs and runs
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if old_text in run.text:
                run.text = run.text.replace(old_text, new_text)
                replaced = True
    
    # If no runs, try the whole text
    if not replaced and old_text in text_frame.text:
        text_frame.text = text_frame.text.replace(old_text, new_text)
        replaced = True
    
    return replaced


def find_and_replace_in_slide(slide, old_text, new_text):
    """
    Find and replace text in all shapes in a slide.
    
    Args:
        slide: The slide object
        old_text: Text to find
        new_text: Text to replace with
        
    Returns:
        int: Number of replacements made
    """
    count = 0
    for shape in slide.shapes:
        if find_and_replace_text_in_shape(shape, old_text, new_text):
            count += 1
    return count


def edit_ompo_report(input_file, output_file, name=None, date=None, status_text=None):
    """
    Edit the OMPO Report PowerPoint presentation.
    
    Args:
        input_file: Path to input PowerPoint file
        output_file: Path to output PowerPoint file
        name: Name of person filling/submitting the report
        date: Date in YYYY-MM-DD format
        status_text: Text to replace in the Status of Burnin Center section
        
    Returns:
        dict: Summary of replacements made
    """
    # Load presentation
    prs = Presentation(input_file)
    
    summary = {
        'name': 0,
        'date': 0,
        'status': 0
    }
    
    # Process each slide
    for slide_idx, slide in enumerate(prs.slides):
        print(f"Processing slide {slide_idx + 1}...")
        
        # Replace Name field
        if name:
            count = find_and_replace_in_slide(slide, "Name", name)
            summary['name'] += count
            if count > 0:
                print(f"  - Replaced 'Name' with '{name}' ({count} occurrences)")
        
        # Replace Date field
        if date:
            count = find_and_replace_in_slide(slide, "YYYY-MM-DD", date)
            summary['date'] += count
            if count > 0:
                print(f"  - Replaced 'YYYY-MM-DD' with '{date}' ({count} occurrences)")
        
        # Replace Status text (the "Text" placeholder in yellow box)
        if status_text:
            count = find_and_replace_in_slide(slide, "Text", status_text)
            summary['status'] += count
            if count > 0:
                print(f"  - Replaced 'Text' with status information ({count} occurrences)")
    
    # Save modified presentation
    prs.save(output_file)
    print(f"\nSaved modified presentation to: {output_file}")
    
    return summary


def main():
    """Main function for command-line usage."""
    parser = argparse.ArgumentParser(
        description='Edit OMPO Report PowerPoint presentation fields',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  %(prog)s input.pptx output.pptx --name "John Doe" --date "2025-11-18"
  %(prog)s template.pptx filled.pptx --name "Jane Smith" --date "2025-11-20" --status "All systems operational. No downtime expected."
  %(prog)s input.pptx output.pptx --date $(date +%%Y-%%m-%%d)  # Use current date
        """
    )
    
    parser.add_argument(
        'input_file',
        type=str,
        help='Input PowerPoint file path'
    )
    
    parser.add_argument(
        'output_file',
        type=str,
        help='Output PowerPoint file path'
    )
    
    parser.add_argument(
        '--name',
        type=str,
        help='Name of person filling/submitting the report'
    )
    
    parser.add_argument(
        '--date',
        type=str,
        help='Date in YYYY-MM-DD format (default: today)'
    )
    
    parser.add_argument(
        '--status',
        type=str,
        help='Status text for the Burnin Center section'
    )
    
    parser.add_argument(
        '--use-today',
        action='store_true',
        help='Use today\'s date automatically'
    )
    
    args = parser.parse_args()
    
    # Validate input file exists
    if not Path(args.input_file).exists():
        print(f"Error: Input file '{args.input_file}' not found.", file=sys.stderr)
        return 1
    
    # Use today's date if requested
    date = args.date
    if args.use_today or (date is None and args.name is None and args.status is None):
        # If no arguments provided, at least use today's date
        date = datetime.now().strftime('%Y-%m-%d')
        print(f"Using today's date: {date}")
    
    # Validate date format if provided
    if date:
        try:
            datetime.strptime(date, '%Y-%m-%d')
        except ValueError:
            print(f"Error: Date must be in YYYY-MM-DD format, got '{date}'", file=sys.stderr)
            return 1
    
    # Check if at least one field is being updated
    if not any([args.name, date, args.status]):
        print("Warning: No fields specified to update. Use --name, --date, or --status options.")
        parser.print_help()
        return 1
    
    try:
        # Edit the presentation
        summary = edit_ompo_report(
            args.input_file,
            args.output_file,
            name=args.name,
            date=date,
            status_text=args.status
        )
        
        print("\nSummary:")
        print(f"  - Name field: {summary['name']} replacements")
        print(f"  - Date field: {summary['date']} replacements")
        print(f"  - Status field: {summary['status']} replacements")
        
        return 0
    
    except Exception as e:
        print(f"Error: {e}", file=sys.stderr)
        import traceback
        traceback.print_exc()
        return 1


if __name__ == '__main__':
    sys.exit(main())
